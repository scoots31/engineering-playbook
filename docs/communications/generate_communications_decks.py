#!/usr/bin/env python3
"""
Generate PowerPoint versions of deck-business.html and deck-solo.html.
Run from repo root or this directory:
  docs/communications/.pptx-build-venv/bin/python generate_communications_decks.py

Outputs:
  deck-business.pptx
  deck-solo.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# ── Design tokens (match deck HTML :root) ─────────────────────────────────
NAVY = RGBColor(0x0F, 0x17, 0x29)
NAVY_MID = RGBColor(0x1A, 0x27, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x2D, 0xD4, 0xBF)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
GOLD_LIGHT = RGBColor(0xFC, 0xD3, 0x4D)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF1, 0xF5, 0xF9)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
SPLIT_BLUE = RGBColor(0x1E, 0x3A, 0x6E)
SPLIT_TEAL = RGBColor(0x0A, 0x3D, 0x38)
SPLIT_GOLD = RGBColor(0x3D, 0x2A, 0x00)
RED_SOFT = RGBColor(0xF8, 0x71, 0x71)
NAVY_TEXT = RGBColor(0x0F, 0x17, 0x29)
SLATE_BODY = RGBColor(0x37, 0x41, 0x51)


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[6]


def _solid_rect(slide, left, top, width, height, rgb: RGBColor):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    sh.line.fill.background()
    return sh


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    *,
    align=PP_ALIGN.LEFT,
    v_anchor=MSO_ANCHOR.TOP,
    word_wrap=True,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = v_anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    return tf


def _para(tf, *, level=0, space_after=Pt(6)):
    if tf.paragraphs and tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = space_after
    return p


def _run(p, text, *, size=Pt(14), bold=False, color=None, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = size
    r.font.bold = bold
    if color is not None:
        r.font.color.rgb = color
    return r


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def build_business_deck(out_path: Path) -> None:
    prs = _new_prs()
    blank = _blank_layout(prs)
    W, H = prs.slide_width, prs.slide_height

    # ── 1 Title ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    m = Inches(0.75)
    tf = _textbox(s, m, Inches(0.85), Inches(11.5), Inches(5.5))
    _run(_para(tf, space_after=Pt(10)), "AI BUILD MINDSET", size=Pt(12), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(14))
    _run(p, "From How We Work ", size=Pt(44), bold=True, color=WHITE)
    _run(p, "Today", size=Pt(44), bold=True, color=BLUE_LIGHT)
    _run(p, " to How We Build ", size=Pt(44), bold=True, color=WHITE)
    _run(p, "Tomorrow", size=Pt(44), bold=True, color=BLUE_LIGHT)
    p2 = _para(tf, space_after=Pt(18))
    _run(
        p2,
        "A framework approach for enabling the AI Build Mindset — built on the discipline that already makes our teams effective.",
        size=Pt(16),
        color=GRAY,
    )
    _solid_rect(s, m, Inches(4.55), Inches(0.65), Inches(0.04), BLUE)

    # ── 2 Roles ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.55), Inches(11.8), Inches(1.1))
    _run(_para(tf), "WHERE WE START", size=Pt(10), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(16))
    _run(p, "Great products come from great team discipline", size=Pt(32), bold=True, color=WHITE)

    cards = [
        ("Product Manager", BLUE_LIGHT, "Defines the problem. Challenges assumptions. Holds the \"why\" before a line of code is written."),
        ("Tech Lead", TEAL_LIGHT, "Validates buildability. Sequences work by risk. Prevents the expensive wrong turns."),
        ("Designer", GOLD_LIGHT, "Produces the shared artifact everyone points at. Surfaces data, flow, and scope questions before build."),
        ("QA", PURPLE_LIGHT, "Enforces the quality gate. Nothing ships without passing the bar — no exceptions."),
    ]
    gx, gy = m, Inches(1.85)
    cw, ch = Inches(5.75), Inches(1.55)
    gap = Inches(0.35)
    for i, (title, tc, desc) in enumerate(cards):
        col, row = i % 2, i // 2
        x = gx + col * (cw + gap)
        y = gy + row * (ch + gap)
        box = _solid_rect(s, x, y, cw, ch, RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = RGBColor(0x33, 0x3D, 0x52)
        ctf = _textbox(s, x + Inches(0.2), y + Inches(0.15), cw - Inches(0.4), ch - Inches(0.25))
        _run(_para(ctf), title, size=Pt(10), bold=True, color=tc)
        _run(_para(ctf, space_after=Pt(4)), desc, size=Pt(11), color=GRAY)

    foot = _textbox(s, m, Inches(5.15), Inches(11.2), Inches(1.35))
    _run(
        _para(foot),
        "This team dynamic is what makes the difference between good product and shipping something that needs to be rebuilt. The discipline is the team.",
        size=Pt(12),
        color=GRAY,
    )

    # ── 3 Split shift ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    split = int(W * 0.54)
    _solid_rect(s, split, 0, W - split, H, SPLIT_BLUE)
    tf = _textbox(s, m, Inches(1.0), split - Inches(1.1), Inches(5.2))
    _run(_para(tf), "THE SHIFT", size=Pt(10), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(12))
    _run(p, "An AI Build Mindset is coming — the question is how it lands", size=Pt(28), bold=True, color=WHITE)
    _run(
        _para(tf),
        "AI tools are enabling individuals to build what previously required a full team. This is real, it is accelerating, and it creates a genuine opportunity for organizations that get the approach right.",
        size=Pt(13),
        color=GRAY,
    )
    _solid_rect(s, split + Inches(0.55), Inches(0.95), Inches(0.45), Inches(0.04), BLUE)
    rtf = _textbox(s, split + Inches(0.55), Inches(1.25), W - split - Inches(1.0), Inches(5.8))
    bullets = [
        "Speed increases significantly — cycles that took weeks can take days",
        "The builder profile changes — domain experts can build without dedicated developers",
        "Output quality varies widely — without structure, results are inconsistent",
        "The discipline gap is the real risk — not the technology itself",
    ]
    for line in bullets:
        p = _para(rtf, space_after=Pt(8))
        _run(p, "→  ", size=Pt(12), bold=True, color=BLUE_LIGHT)
        _run(p, line, size=Pt(12), color=GRAY_LIGHT)

    # ── 4 Gap list ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.55), Inches(11.5), Inches(6.5))
    _run(_para(tf), "THE REAL PROBLEM", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(14)), "When the team goes away, so does the discipline", size=Pt(30), bold=True, color=WHITE)
    items = [
        "No PM challenge — features get built without validating the problem first. Assumptions go unchecked.",
        "No tech review — work gets sequenced by comfort, not by risk. Expensive wrong turns happen late.",
        "No shared design artifact — everyone works from a description. Scope, data gaps, and flow issues surface in code.",
        "No QA gate — \"it works on my machine\" replaces structured verification. Quality is inconsistent.",
        "No process accountability — what was agreed in discovery gets quietly lost during build.",
    ]
    for line in items:
        p = _para(tf, space_after=Pt(6))
        _run(p, "✕  ", size=Pt(12), bold=True, color=RED_SOFT)
        _run(p, line, size=Pt(12), color=GRAY_LIGHT)

    # ── 5 Framework mapping split ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    split = int(W * 0.54)
    _solid_rect(s, split, 0, W - split, H, SPLIT_TEAL)
    tf = _textbox(s, m, Inches(1.0), split - Inches(1.0), Inches(5.0))
    _run(_para(tf), "OUR APPROACH", size=Pt(10), bold=True, color=TEAL_LIGHT)
    p = _para(tf, space_after=Pt(10))
    _run(p, "What if the framework ", size=Pt(28), bold=True, color=WHITE)
    _run(p, "replaced the team?", size=Pt(28), bold=True, color=TEAL_LIGHT)
    _run(
        _para(tf),
        "The roles that make teams effective don't disappear — they get embedded in the framework. Every phase has specialist thinking. Every handoff has a gate.",
        size=Pt(12),
        color=GRAY,
    )
    _solid_rect(s, split + Inches(0.5), Inches(0.95), Inches(0.45), Inches(0.04), TEAL)
    rows = [
        ("PM challenge", "discover · grill-me"),
        ("Tech Lead review", "principal-engineer · prd-to-plan"),
        ("Design review", "design-sprint · design-review"),
        ("QA gate", "code-review · solo-qa · phase-test"),
        ("Process mapping", "process-mapper (always-on)"),
        ("Institutional memory", "product-continuity (always-on)"),
    ]
    rtf = _textbox(s, split + Inches(0.45), Inches(1.2), W - split - Inches(0.9), Inches(5.8))
    hdr = _para(rtf, space_after=Pt(6))
    _run(hdr, "Team discipline → Framework skill", size=Pt(9), bold=True, color=TEAL_LIGHT)
    for left, right in rows:
        p = _para(rtf, space_after=Pt(4))
        _run(p, f"{left} → ", size=Pt(11), color=GRAY)
        _run(p, right, size=Pt(11), bold=True, color=WHITE)

    # ── 6 Process thread ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.6), Inches(6.6))
    _run(_para(tf), "THE BACKBONE", size=Pt(10), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(12))
    _run(p, "Process is the thread that runs through ", size=Pt(30), bold=True, color=WHITE)
    _run(p, "everything", size=Pt(30), bold=True, color=BLUE_LIGHT)
    pts = [
        "Every project starts with a process conversation — how does this work today? What should it look like when we're done?",
        "The to-be map becomes the contract — agreed before design begins, held as the accountability standard through build and testing.",
        "Every feature slice maps to a process step — nothing gets built that can't be traced to an agreed process outcome.",
        "Phase test validates process coverage — did we build what the process requires?",
    ]
    for i, t in enumerate(pts, 1):
        p = _para(tf, space_after=Pt(6))
        _run(p, f"{i}  ", size=Pt(12), bold=True, color=GOLD_LIGHT)
        _run(p, t, size=Pt(12), color=GRAY_LIGHT)
    _run(
        _para(tf, space_after=Pt(0)),
        "This is what your process mapping team does today. The framework brings it into every build, automatically.",
        size=Pt(11),
        color=GRAY,
    )

    # ── 7 Lifecycle + stats ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.8), Inches(2.0))
    _run(_para(tf), "END-TO-END", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(10)), "A complete lifecycle — from first idea to deployed and tested", size=Pt(26), bold=True, color=WHITE)
    flow = "Start › Brainstorm › Discover › Tech Context › Design Sprint › Design Review › Plan › Build › QA Chain › Phase Test › Deploy"
    _run(_para(tf, space_after=Pt(10)), flow, size=Pt(11), color=GRAY_LIGHT)

    sx, sy = m, Inches(3.15)
    stats = [
        ("24", "Specialist skills covering the full lifecycle"),
        ("4", "Always-on skills running without user invocation"),
        ("5", "Hard gates before code ships to production"),
        ("7", "Specialist lenses in the Phase Test alone"),
    ]
    bw = (W - 2 * m - Inches(0.45)) // 4
    for i, (num, lbl) in enumerate(stats):
        x = sx + i * (bw + Inches(0.15))
        box = _solid_rect(s, x, sy, bw, Inches(1.35), RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = RGBColor(0x33, 0x3D, 0x52)
        stf = _textbox(s, x, sy + Inches(0.15), bw, Inches(1.1), align=PP_ALIGN.CENTER)
        stf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(_para(stf), num, size=Pt(36), bold=True, color=WHITE)
        p = _para(stf, space_after=Pt(2))
        p.alignment = PP_ALIGN.CENTER
        _run(p, lbl, size=Pt(10), color=GRAY)

    # ── 8 Light familiar ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, OFF_WHITE)
    tf = _textbox(s, m, Inches(0.55), Inches(11.6), Inches(6.4))
    _run(_para(tf), "FAMILIAR BY DESIGN", size=Pt(10), bold=True, color=BLUE)
    _run(
        _para(tf, space_after=Pt(12)),
        "This isn't new behavior. It's familiar discipline in a new form.",
        size=Pt(28),
        bold=True,
        color=NAVY_TEXT,
    )
    light_bullets = [
        ("Process maps", "your team already does this. The framework formalizes as-is and to-be maps in every project."),
        ("PM, Tech Lead, Designer, QA roles", "these exist in the framework as specialist skills."),
        ("Review gates", "nothing advances without passing the gate."),
        ("Institutional memory", "decisions, assumptions, risks, and changes are captured continuously."),
        ("Continuous improvement", "the retrospective skill feeds improvements back into the framework."),
    ]
    for head, rest in light_bullets:
        p = _para(tf, space_after=Pt(6))
        _run(p, "✓  ", size=Pt(12), bold=True, color=BLUE)
        _run(p, head + " — ", size=Pt(12), bold=True, color=NAVY_TEXT)
        _run(p, rest, size=Pt(12), color=SLATE_BODY)

    # ── 9 Differentiation ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.55), Inches(11.8), Inches(6.6))
    _run(_para(tf), "FACTUAL DIFFERENTIATION", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(12)), "What makes this framework different", size=Pt(30), bold=True, color=WHITE)
    diffs = [
        ("Process-first, not code-first", "Every project starts with a process conversation, not a prompt."),
        ("Full lifecycle coverage", "From idea through deployed and verified — not just a coding assistant."),
        ("Team discipline without a team", "PM, Tech Lead, Designer, QA, and Process Mapper roles are embedded as skills."),
        ("Self-improving", "The retrospective skill captures what works from real usage."),
        ("Tool-agnostic & collaborative", "Built to work with Claude Code and Cursor on the same project."),
    ]
    for label, desc in diffs:
        p = _para(tf, space_after=Pt(6))
        _run(p, label + " — ", size=Pt(12), bold=True, color=BLUE_LIGHT)
        _run(p, desc, size=Pt(12), color=GRAY_LIGHT)

    # ── 10 Pilot ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.6), Inches(1.0))
    _run(_para(tf), "WHERE TO START", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(12)), "A pilot approach — one project, one builder", size=Pt(28), bold=True, color=WHITE)
    steps = [
        ("Step 1", "Select the right first project", "A real project with a real process behind it — not a toy."),
        ("Step 2", "Run Discover and Design Sprint", "Does the framework produce a process map and a visual artifact the team recognizes?"),
        ("Step 3", "Build through the QA chain", "Run build → code review → QA on at least one phase."),
        ("Step 4", "Evaluate and iterate", "Use retrospective output to refine the framework for your context."),
    ]
    pw = (W - 2 * m - Inches(0.35)) / 2
    ph = Inches(1.45)
    for i, (st, ti, de) in enumerate(steps):
        col, row = i % 2, i // 2
        x = m + col * (pw + Inches(0.35))
        y = Inches(1.55) + row * (ph + Inches(0.2))
        box = _solid_rect(s, x, y, pw, ph, RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = RGBColor(0x40, 0x4A, 0x5E)
        ctf = _textbox(s, x + Inches(0.2), y + Inches(0.12), pw - Inches(0.35), ph - Inches(0.2))
        _run(_para(ctf), st.upper(), size=Pt(9), bold=True, color=BLUE_LIGHT)
        _run(_para(ctf), ti, size=Pt(14), bold=True, color=WHITE)
        _run(_para(ctf), de, size=Pt(11), color=GRAY)

    # ── 11 Closing ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY_MID)
    tf = _textbox(s, Inches(1.2), Inches(2.2), Inches(11.0), Inches(3.5), align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = _para(tf, space_after=Pt(14))
    p.alignment = PP_ALIGN.CENTER
    _run(p, "The discipline doesn't have to leave\nwhen the ", size=Pt(36), bold=True, color=WHITE)
    _run(p, "team evolves.", size=Pt(36), bold=True, color=BLUE_LIGHT)
    p2 = _para(tf, space_after=Pt(18))
    p2.alignment = PP_ALIGN.CENTER
    _run(
        p2,
        "We built a framework that carries the judgment, the gates, and the process accountability of a great team into an AI Build Mindset — without rebuilding from scratch what already works.",
        size=Pt(15),
        color=GRAY,
    )
    pills = ["24 Skills", "Full Lifecycle", "Process-First", "Tool-Agnostic"]
    px = Inches(2.4)
    for lab in pills:
        pill = _solid_rect(s, px, Inches(5.35), Inches(1.85), Inches(0.42), RGBColor(0x22, 0x2E, 0x48))
        pill.line.color.rgb = BLUE
        ptf = _textbox(s, px, Inches(5.38), Inches(1.85), Inches(0.4), align=PP_ALIGN.CENTER)
        ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(_para(ptf), lab, size=Pt(11), bold=True, color=BLUE_LIGHT)
        px += Inches(2.05)

    prs.save(out_path)


def build_solo_deck(out_path: Path) -> None:
    prs = _new_prs()
    blank = _blank_layout(prs)
    W, H = prs.slide_width, prs.slide_height
    m = Inches(0.75)

    # 1 Title
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.85), Inches(11.5), Inches(5.5))
    _run(_para(tf, space_after=Pt(10)), "SOLO BUILDER FRAMEWORK", size=Pt(12), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(14))
    _run(p, "Team discipline.\n", size=Pt(44), bold=True, color=WHITE)
    _run(p, "Without the team.", size=Pt(44), bold=True, color=BLUE_LIGHT)
    _run(
        _para(tf),
        "A complete lifecycle framework for building software solo with AI — from first idea through deployed and tested.",
        size=Pt(16),
        color=GRAY,
    )
    _solid_rect(s, m, Inches(4.55), Inches(0.65), Inches(0.04), BLUE)

    # 2 Problem
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.55), Inches(11.8), Inches(6.5))
    _run(_para(tf), "THE PROBLEM", size=Pt(10), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(12))
    _run(p, "Solo builders using AI ", size=Pt(30), bold=True, color=WHITE)
    _run(p, "drift toward waterfall", size=Pt(30), bold=True, color=GOLD_LIGHT)
    probs = [
        "No external forcing function — no standup, no sprint review, no teammate waiting.",
        "AI defaults to building what you describe — not challenging whether you're ready to build it.",
        "No shared artifact before code starts — scope, data gaps, and flow issues surface in code.",
        "Quality gates disappear — \"it works\" becomes the bar. What ships is inconsistent.",
    ]
    for line in probs:
        p = _para(tf, space_after=Pt(6))
        _run(p, "✕  ", size=Pt(12), bold=True, color=RED_SOFT)
        _run(p, line, size=Pt(12), color=GRAY_LIGHT)
    _run(
        _para(tf),
        "Replace \"phase complete\" with \"enough to learn from.\" The framework provides the forcing functions a team would provide.",
        size=Pt(11),
        color=GRAY,
    )

    # 3 Entry points
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.55), Inches(11.6), Inches(1.05))
    _run(_para(tf), "WHERE YOU START", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(12)), "Two entry points — the framework routes for you", size=Pt(28), bold=True, color=WHITE)
    ew = (W - 2 * m - Inches(0.35)) / 2
    eh = Inches(2.35)
    entries = [
        ("ENTRY 1 — You know what you want to build", "The idea is clear. Skip brainstorming — go straight to Discover.", "Signals: specific solution, clear user, defined outcome"),
        ("ENTRY 2 — You're working through an idea", "Brainstorming gives the idea something concrete to react to.", "Signals: exploratory tone, uncertain scope or user"),
    ]
    for i, (t, d, sig) in enumerate(entries):
        x = m + i * (ew + Inches(0.35))
        y = Inches(1.45)
        box = _solid_rect(s, x, y, ew, eh, RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = BLUE if i == 0 else GOLD
        ctf = _textbox(s, x + Inches(0.22), y + Inches(0.18), ew - Inches(0.4), eh - Inches(0.3))
        _run(_para(ctf), t, size=Pt(11), bold=True, color=BLUE_LIGHT if i == 0 else GOLD_LIGHT)
        _run(_para(ctf), d, size=Pt(12), color=GRAY)
        _run(_para(ctf), sig, size=Pt(10), color=GRAY_LIGHT)
    foot = _textbox(s, m, Inches(4.05), Inches(11.5), Inches(0.9))
    _run(
        _para(foot),
        "The start skill reads your opening message and routes automatically — both paths converge at Discover.",
        size=Pt(11),
        color=GRAY,
    )

    # 4 Lifecycle overview
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.8), Inches(6.8))
    _run(_para(tf), "THE FULL LIFECYCLE", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(10)), "Every phase has a purpose, an output, and a gate", size=Pt(28), bold=True, color=WHITE)
    _run(
        _para(tf, space_after=Pt(10)),
        "Brainstorm › Discover › Tech Context › Design Sprint › Design Review › Plan › Build › QA Chain › Phase Test › Deploy",
        size=Pt(11),
        color=GRAY_LIGHT,
    )
    trio = [
        ("You bring:", "the idea, domain knowledge, warmer/colder judgment, and final decisions", BLUE_LIGHT),
        ("The framework brings:", "specialist thinking, design execution, quality gates, process accountability, memory", TEAL_LIGHT),
        ("Gates stop advancement", "until conditions are met — to-be map agreed, Ready slices, phase test OPEN before deploy", GOLD_LIGHT),
    ]
    for head, body, hc in trio:
        p = _para(tf, space_after=Pt(6))
        _run(p, "→  ", size=Pt(12), bold=True, color=hc)
        _run(p, head + " ", size=Pt(12), bold=True, color=WHITE)
        _run(p, body, size=Pt(12), color=GRAY_LIGHT)

    # 5 Always-on
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.6), Inches(1.1))
    _run(_para(tf), "ALWAYS-ON", size=Pt(10), bold=True, color=GOLD_LIGHT)
    p = _para(tf, space_after=Pt(10))
    _run(p, "Four skills that run the entire time — ", size=Pt(26), bold=True, color=WHITE)
    _run(p, "never invoked by you", size=Pt(26), bold=True, color=GOLD_LIGHT)
    skills = [
        ("process-mapper", "Process accountability — as-is / to-be maps; to-be is the contract.", BLUE_LIGHT),
        ("product-continuity", "Institutional memory — decisions, risks, glossary across sessions.", TEAL_LIGHT),
        ("framework-health", "Background monitor — surfaces one issue at a time; never blocks.", GOLD_LIGHT),
        ("retrospective", "Continuous learning — flags in the moment, processes at phase end.", PURPLE_LIGHT),
    ]
    gw = (W - 2 * m - Inches(0.35)) / 2
    gh = Inches(1.55)
    for i, (name, desc, ac) in enumerate(skills):
        col, row = i % 2, i // 2
        x = m + col * (gw + Inches(0.35))
        y = Inches(1.45) + row * (gh + Inches(0.22))
        box = _solid_rect(s, x, y, gw, gh, RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = RGBColor(0x40, 0x4A, 0x5E)
        ctf = _textbox(s, x + Inches(0.2), y + Inches(0.12), gw - Inches(0.35), gh - Inches(0.2))
        _run(_para(ctf), name, size=Pt(13), bold=True, color=ac)
        _run(_para(ctf), desc, size=Pt(11), color=GRAY)

    # 6 Split design artifact
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    split = int(W * 0.54)
    _solid_rect(s, split, 0, W - split, H, SPLIT_BLUE)
    tf = _textbox(s, m, Inches(0.95), split - Inches(1.0), Inches(5.0))
    _run(_para(tf), "PHASE 2", size=Pt(10), bold=True, color=BLUE_LIGHT)
    p = _para(tf, space_after=Pt(8))
    _run(p, "The design artifact is the ", size=Pt(26), bold=True, color=WHITE)
    _run(p, "Rosetta Stone", size=Pt(26), bold=True, color=BLUE_LIGHT)
    _run(
        _para(tf),
        "Design Sprint produces an HTML artifact — interactive, good enough to decide from. It surfaces data questions, process gaps, scope decisions, and implied infrastructure.",
        size=Pt(12),
        color=GRAY,
    )
    _solid_rect(s, split + Inches(0.55), Inches(0.95), Inches(0.45), Inches(0.04), BLUE)
    rights = [
        "Data questions — where does each field come from?",
        "Process gaps — missing steps and handoffs surface in flow",
        "Scope decisions — every element is build/defer explicitly",
        "Implied infrastructure — toggles may hide pipelines",
    ]
    rtf = _textbox(s, split + Inches(0.5), Inches(1.2), W - split - Inches(0.95), Inches(5.8))
    for line in rights:
        p = _para(rtf, space_after=Pt(8))
        _run(p, "•  ", size=Pt(12), bold=True, color=BLUE_LIGHT)
        _run(p, line, size=Pt(12), color=GRAY_LIGHT)

    # 7 Four anchors
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.6), Inches(1.0))
    _run(_para(tf), "PHASE 4 — BUILD", size=Pt(10), bold=True, color=TEAL_LIGHT)
    p = _para(tf, space_after=Pt(10))
    _run(p, "No slice starts without ", size=Pt(28), bold=True, color=WHITE)
    _run(p, "four anchors", size=Pt(28), bold=True, color=TEAL_LIGHT)
    anchors = [
        ("Design Anchor", BLUE_LIGHT, "Exact screen/element from the design sprint."),
        ("Data Anchor", TEAL_LIGHT, "Mock data fields; data-scaffold generates realistic values."),
        ("Done Anchor", GOLD_LIGHT, "2–3 verifiable criteria — specific, checkable statements."),
        ("Process Anchor", PURPLE_LIGHT, "Which to-be process step this slice implements."),
    ]
    aw = (W - 2 * m - Inches(0.35)) / 2
    ah = Inches(1.38)
    for i, (title, colc, desc) in enumerate(anchors):
        col, row = i % 2, i // 2
        x = m + col * (aw + Inches(0.35))
        y = Inches(1.45) + row * (ah + Inches(0.2))
        box = _solid_rect(s, x, y, aw, ah, RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = RGBColor(0x40, 0x4A, 0x5E)
        ctf = _textbox(s, x + Inches(0.2), y + Inches(0.12), aw - Inches(0.35), ah - Inches(0.2))
        _run(_para(ctf), title, size=Pt(11), bold=True, color=colc)
        _run(_para(ctf), desc, size=Pt(11), color=GRAY)
    foot = _textbox(s, m, Inches(5.05), Inches(11.2), Inches(0.85))
    _run(
        _para(foot),
        "If a slice has no process anchor, the framework stops and asks whether you're building outside the agreed process.",
        size=Pt(10),
        color=GRAY,
    )

    # 8 QA chain
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.8), Inches(2.0))
    _run(_para(tf), "AUTOMATIC QA", size=Pt(10), bold=True, color=TEAL_LIGHT)
    p = _para(tf, space_after=Pt(10))
    _run(p, "Code-complete triggers a ", size=Pt(26), bold=True, color=WHITE)
    _run(p, "three-skill QA chain", size=Pt(26), bold=True, color=TEAL_LIGHT)
    chain = [
        ("code-review-and-quality", "7 checks on patterns, data sourcing, scope, documentation — pass invokes solo-qa"),
        ("solo-qa Part 1", "Active testing with evidence; cross-reference UI to mock data"),
        ("solo-qa Part 2", "You sign off in the browser — visual, behavioral, done criteria"),
    ]
    cw = (W - 2 * m - Inches(0.5)) / 3
    cy = Inches(1.45)
    for i, (name, desc) in enumerate(chain):
        x = m + i * (cw + Inches(0.25))
        box = _solid_rect(s, x, cy, cw, Inches(2.55), RGBColor(0x18, 0x22, 0x38))
        box.line.color.rgb = RGBColor(0x40, 0x4A, 0x5E)
        ctf = _textbox(s, x + Inches(0.15), cy + Inches(0.12), cw - Inches(0.28), Inches(2.35))
        _run(_para(ctf), name, size=Pt(12), bold=True, color=WHITE)
        _run(_para(ctf), desc, size=Pt(10), color=GRAY)
        if i < 2:
            arr = _textbox(s, x + cw + Inches(0.02), cy + Inches(1.05), Inches(0.22), Inches(0.5))
            _run(_para(arr), "›", size=Pt(22), bold=True, color=BLUE_LIGHT)
    foot = _textbox(s, m, Inches(4.25), Inches(11.5), Inches(0.85))
    _run(
        _para(foot),
        "qa-triage classifies unexpected findings: bug / missing requirement / regression — then routes correctly.",
        size=Pt(10),
        color=GRAY,
    )

    # 9 Process thread split
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    split = int(W * 0.54)
    _solid_rect(s, split, 0, W - split, H, SPLIT_GOLD)
    tf = _textbox(s, m, Inches(1.0), split - Inches(1.0), Inches(5.0))
    _run(_para(tf), "THE BACKBONE", size=Pt(10), bold=True, color=GOLD_LIGHT)
    p = _para(tf, space_after=Pt(10))
    _run(p, "The process thread runs from ", size=Pt(26), bold=True, color=WHITE)
    _run(p, "first conversation to final test", size=Pt(26), bold=True, color=GOLD_LIGHT)
    _run(
        _para(tf),
        "The to-be process map is the contract every downstream phase is held to — automatically.",
        size=Pt(12),
        color=GRAY,
    )
    _solid_rect(s, split + Inches(0.55), Inches(0.95), Inches(0.45), Inches(0.04), GOLD)
    steps = [
        ("Discover", "As-is and to-be maps agreed. The to-be map is the contract."),
        ("Design Sprint", "Screens cross-referenced to the map; gaps become decisions."),
        ("Design Review", "Coverage each round; every slice gets a process anchor before Ready."),
        ("Build", "Process anchor required; stops if outside agreed process."),
        ("Phase Test", "Use cases from the map — did we build what the process requires?"),
    ]
    rtf = _textbox(s, split + Inches(0.5), Inches(1.15), W - split - Inches(0.95), Inches(5.9))
    for i, (ph, tx) in enumerate(steps, 1):
        p = _para(rtf, space_after=Pt(6))
        _run(p, f"{i}  ", size=Pt(11), bold=True, color=GOLD_LIGHT)
        _run(p, f"{ph}: ", size=Pt(11), bold=True, color=WHITE)
        _run(p, tx, size=Pt(11), color=GRAY_LIGHT)

    # 10 Phase test
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.5), Inches(11.8), Inches(6.8))
    _run(_para(tf), "PHASE 5 — EXPLICIT", size=Pt(10), bold=True, color=PURPLE_LIGHT)
    p = _para(tf, space_after=Pt(10))
    _run(p, "Phase Test: ", size=Pt(28), bold=True, color=WHITE)
    _run(p, "seven specialist lenses", size=Pt(28), bold=True, color=PURPLE_LIGHT)
    stages = [
        "Stage 1 — Environment readiness (mock inactive, slices deployed, clean load)",
        "Stage 2 — Use case creator from discovery + to-be map",
        "Stage 3 — Data specialist (mock→real, API verification)",
        "Stage 4+5 — Tester + regression with named evidence",
        "Stage 6 — Acceptance reviewer (PM lens on intent)",
        "Stage 7 — Gate decision: OPEN or HOLD with targeted re-test",
    ]
    for line in stages:
        p = _para(tf, space_after=Pt(4))
        _run(p, line, size=Pt(11), color=GRAY_LIGHT)
    _run(
        _para(tf),
        "Invoke /phase-test explicitly when the phase is built — deliberate, not automatic.",
        size=Pt(10),
        color=GRAY,
    )

    # 11 Differentiation
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY)
    tf = _textbox(s, m, Inches(0.55), Inches(11.8), Inches(6.6))
    _run(_para(tf), "FACTUAL DIFFERENTIATION", size=Pt(10), bold=True, color=BLUE_LIGHT)
    _run(_para(tf, space_after=Pt(12)), "What makes this different", size=Pt(30), bold=True, color=WHITE)
    diffs = [
        ("Not a prompt collection", "24 interconnected skills; QA chain runs without manual triggers."),
        ("Process contract", "To-be map agreed before design; held through Phase Test."),
        ("Design before code", "HTML artifact before build; scope/data/flow surface early."),
        ("Quality with evidence", "solo-qa cross-references UI to mock data — named evidence."),
        ("Self-improving", "Retrospective captures learnings; playbook improvements queue from real use."),
    ]
    for label, desc in diffs:
        p = _para(tf, space_after=Pt(6))
        _run(p, label + " — ", size=Pt(12), bold=True, color=BLUE_LIGHT)
        _run(p, desc, size=Pt(12), color=GRAY_LIGHT)

    # 12 Closing
    s = prs.slides.add_slide(blank)
    _solid_rect(s, 0, 0, W, H, NAVY_MID)
    tf = _textbox(s, Inches(1.2), Inches(2.15), Inches(11.0), Inches(3.6), align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = _para(tf, space_after=Pt(14))
    p.alignment = PP_ALIGN.CENTER
    _run(p, "Start with one idea.\nRun it ", size=Pt(36), bold=True, color=WHITE)
    _run(p, "end to end.", size=Pt(36), bold=True, color=TEAL_LIGHT)
    p2 = _para(tf, space_after=Pt(18))
    p2.alignment = PP_ALIGN.CENTER
    _run(
        p2,
        "Open a new project. Let the framework route you — Brainstorm or Discover. Stay through Design Sprint, Build, and Phase Test.",
        size=Pt(15),
        color=GRAY,
    )
    pills = ["24 Skills", "Automatic QA Chain", "Process-First", "Claude Code + Cursor"]
    px = Inches(1.85)
    for lab in pills:
        pill = _solid_rect(s, px, Inches(5.35), Inches(2.05), Inches(0.42), RGBColor(0x22, 0x2E, 0x48))
        pill.line.color.rgb = TEAL
        ptf = _textbox(s, px, Inches(5.38), Inches(2.05), Inches(0.4), align=PP_ALIGN.CENTER)
        ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(_para(ptf), lab, size=Pt(10), bold=True, color=TEAL_LIGHT)
        px += Inches(2.22)

    prs.save(out_path)


def main():
    base = Path(__file__).resolve().parent
    build_business_deck(base / "deck-business.pptx")
    build_solo_deck(base / "deck-solo.pptx")
    print(f"Wrote:\n  {base / 'deck-business.pptx'}\n  {base / 'deck-solo.pptx'}")


if __name__ == "__main__":
    main()
